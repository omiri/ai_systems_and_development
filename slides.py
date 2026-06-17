from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)  # Widescreen 16:9
prs.slide_height = Inches(7.5)

# --- SLIDE 1: Title Slide ---
slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "AI SYSTEMS & DEPLOYMENT"
slide.placeholders[1].text = "AnalystLab Africa | Batch A Portfolio | Week 7\nPresenter: Pius Omirida Otumala"

# --- SLIDE 2: Journey to Production ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "The Journey to Production"
tf = slide.placeholders[1].text_frame
tf.text = "Transitioning from isolated notebooks to live software systems."
tf.add_paragraph().text = "• Shifting gears from model training to production inference utilities."
tf.add_paragraph().text = "• Bypassing local runtime execution boundaries to build web accessible apps."

# --- SLIDE 4: Task 1 ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Task 1: Model Serialization"
tf = slide.placeholders[1].text_frame
tf.text = "Freezing the Intelligence using Joblib:"
tf.add_paragraph().text = "• Captured weights and decision parameters into a clean '.joblib' binary layout."
tf.add_paragraph().text = "• Completely avoids the runtime cost of retraining the model on incoming data."

# --- SLIDE 5: Task 2 ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Task 2: API Architecture"
tf = slide.placeholders[1].text_frame
tf.text = "The Serving Layer with FastAPI:"
tf.add_paragraph().text = "• Built asynchronous endpoint routes using Python framework rules."
tf.add_paragraph().text = "• Enforced clear model inputs via Pydantic type validation constraints."

# --- SLIDE 6: Task 3 ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Task 3: Cloud Deployment"
tf = slide.placeholders[1].text_frame
tf.text = "Hosting on Production Infrastructure:"
tf.add_paragraph().text = "• Version controlled via GitHub to ensure smooth source tracking mappings."
tf.add_paragraph().text = "• Hosted on Render Web Services to deliver a live, public interface URL endpoint."

# --- SLIDE 7: Task 4 ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Task 4: Validation Success"
tf = slide.placeholders[1].text_frame
tf.text = "Verification and System Metrics:"
tf.add_paragraph().text = "• Confirmed 82% baseline accuracy mapping trends on the model core."
tf.add_paragraph().text = "• Verified live JSON responses show ~81.5% survival probability for first-class female inputs."

# --- SLIDE 12: Conclusion ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "QUESTIONS?"
slide.placeholders[1].text = "Thank you!\nGitHub: github.com/omiri/ai_systems_and_development\nPresenter: Pius Omirida Otumala"

prs.save("AI_Deployment_Portfolio.pptx")
print("PowerPoint file generated successfully as AI_Deployment_Portfolio.pptx!")